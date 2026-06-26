#!/usr/bin/env python3
"""
color_unify.py — 混合方案色系统一后处理工具
用途：将工作流C咨询模板的默认配色替换为用户选定的设计哲学色板
用法：python3 color_unify.py input.pptx output.pptx --bg 0D0D2B --primary 00FF88 --accent FF6B35

内置默认替换映射（工作流C标准默认色 → 自定义色）：
  #2E9BD6 → primary（主色）
  #0A1F3D → bg（背景）
  #1F3D6B → bg
  #003366 → bg
  #E04E5E → accent（强调色）
  #4CAF50 → primary
  #F2F2F2 → card_bg（卡片背景，默认bg加深20%）
  #333333 → text（文字色，默认CCCCFF）
  #666666 → text_muted（暗文字，默认AAAACC）
"""
import argparse
from pptx import Presentation
from pptx.dml.color import RGBColor


def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def replace_colors(elem, theme_map: dict):
    """递归替换 XML 元素中的 srgbClr 颜色值"""
    for child in elem.iter():
        if child.tag.endswith('}srgbClr'):
            val = child.get('val', '').upper()
            if val in theme_map:
                child.set('val', theme_map[val])


def build_theme_map(primary: str, accent: str, bg: str, text: str = "E0E0FF",
                    text_muted: str = "AAAACC", card_bg: str = None) -> dict:
    if card_bg is None:
        # 卡片背景默认比主背景亮一档
        r, g, b = int(bg[0:2], 16), int(bg[2:4], 16), int(bg[4:6], 16)
        card_bg = f"{min(r+30, 255):02X}{min(g+30, 255):02X}{min(b+30, 255):02X}"
    return {
        "2E9BD6": primary,
        "0A1F3D": bg,
        "1F3D6B": bg,
        "003366": bg,
        "0D1B2A": bg,
        "E04E5E": accent,
        "4CAF50": primary,
        "F2F2F2": card_bg,
        "F0F0F0": card_bg,
        "333333": text,
        "444444": text,
        "666666": text_muted,
        "C0C0C0": text_muted,
        "F4C57A": "FFD700",   # 黄色点缀保留为金色
    }


def main():
    parser = argparse.ArgumentParser(description="统一 PPTX 混合方案色系")
    parser.add_argument("input", help="输入 PPTX 路径")
    parser.add_argument("output", help="输出 PPTX 路径")
    parser.add_argument("--bg", required=True, help="背景色 hex（不含#），如 0D0D2B")
    parser.add_argument("--primary", required=True, help="主色 hex，如 00FF88")
    parser.add_argument("--accent", required=True, help="强调色 hex，如 FF6B35")
    parser.add_argument("--text", default="E0E0FF", help="正文色 hex，默认 E0E0FF")
    parser.add_argument("--card-bg", default=None, help="卡片背景色 hex，默认自动生成")
    parser.add_argument("--dark-bg", action="store_true", default=True,
                        help="将所有页面背景改为 --bg 色（默认开启）")
    args = parser.parse_args()

    theme_map = build_theme_map(
        primary=args.primary.upper(),
        accent=args.accent.upper(),
        bg=args.bg.upper(),
        text=args.text.upper(),
        card_bg=args.card_bg.upper() if args.card_bg else None,
    )
    bg_color = hex_to_rgb(args.bg)

    prs = Presentation(args.input)
    for slide in prs.slides:
        if args.dark_bg:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = bg_color
        for shape in slide.shapes:
            replace_colors(shape._element, theme_map)

    prs.save(args.output)
    print(f"✅ 色系统一完成: {args.output}")
    print(f"   bg={args.bg}  primary={args.primary}  accent={args.accent}")


if __name__ == "__main__":
    main()
